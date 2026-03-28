import json
import os
import sys

from pptx import Presentation

sys.path.insert(0, os.path.abspath(os.path.join(os.path.dirname(__file__), "../..")))

from phemacast.castrs.pptx_castr import PPTXCastr
from prompits.pools.filesystem import FileSystemPool
from prompits.tests.test_support import build_agent_from_config


def _collect_slide_text(slide) -> str:
    return "\n".join(
        shape.text.strip()
        for shape in slide.shapes
        if hasattr(shape, "text") and shape.text and shape.text.strip()
    )


def test_pptx_castr_generates_real_powerpoint_with_expected_content(tmp_path):
    pool = FileSystemPool("ppt_pool", "Temporary storage for pptx castr tests", str(tmp_path))
    castr = PPTXCastr(name="PptCastr", pool=pool, auto_register=False)

    phema = {
        "name": "Weekly Market Brief",
        "description": "A concise cross-asset update for leadership.",
        "owner": "MacroPhemar",
        "tags": ["macro", "markets"],
        "sections": [
            {
                "name": "Macro Outlook",
                "description": "Key backdrop for the week ahead.",
                "content": [
                    "Inflation cooled for a second consecutive reading.",
                    {"pulse_name": "rates", "result": {"value": "Treasury yields eased across the curve."}},
                ],
            },
            {
                "name": "Equities",
                "content": [
                    {"key": "sp500", "data": {"close": 5123.45, "change_pct": 1.2}},
                ],
            },
        ],
    }

    result = castr.cast(
        phema,
        format="ppt",
        preferences={"audience": "Executive team", "language": "en", "theme": "Midnight Executive"},
    )

    assert result["status"] == "success"
    assert result["format"] == "PPTX"
    assert result["location"].startswith("media/")
    assert result["location"].endswith(".pptx")

    output_path = tmp_path / result["location"]
    assert output_path.exists()

    presentation = Presentation(output_path)
    assert len(presentation.slides) == 4

    deck_text = "\n".join(_collect_slide_text(slide) for slide in presentation.slides)
    assert "Weekly Market Brief" in deck_text
    assert "Executive team" in deck_text
    assert "Macro Outlook" in deck_text
    assert "Inflation cooled for a second consecutive reading." in deck_text
    assert "Treasury yields eased across the curve." in deck_text
    assert "5123.45" in deck_text


def test_build_agent_from_config_loads_specialized_pptx_castr(tmp_path):
    storage_dir = tmp_path / "storage"
    config_path = tmp_path / "ppt.castr"
    config_path.write_text(
        json.dumps(
            {
                "name": "PptCastr",
                "type": "phemacast.castrs.pptx_castr.PptCastr",
                "host": "127.0.0.1",
                "port": 8038,
                "plaza_url": None,
                "role": "castr",
                "tags": ["castr", "pptx", "media", "render"],
                "pools": [
                    {
                        "type": "FileSystemPool",
                        "name": "ppt_castr_pool",
                        "description": "Local pool for PPT Castr state.",
                        "root_path": str(storage_dir),
                    }
                ],
                "castr": {
                    "description": "A Castr agent specialized in rendering Phemas into PPTX presentations.",
                    "tags": ["pptx"],
                    "media_type": "PPTX",
                },
            }
        ),
        encoding="utf-8",
    )

    agent = build_agent_from_config(str(config_path))

    assert isinstance(agent, PPTXCastr)
    assert agent.media_type == "PPTX"

    result = agent.cast({"name": "Config Loaded Deck", "sections": []}, format="pptx")
    assert (storage_dir / result["location"]).exists()
