import csv
import json
import os
import sys
from pathlib import Path

from fpdf import FPDF
from pptx import Presentation


sys.path.insert(0, os.path.abspath(os.path.join(os.path.dirname(__file__), "../..")))

from attas.utils.report_to_phema import (
    DEFAULT_OLLAMA_KEEP_ALIVE,
    DEFAULT_OLLAMA_MODELS,
    DEFAULT_OLLAMA_TIMEOUT,
    DEFAULT_OLLAMA_URL,
    convert_report_file,
    convert_report_payload,
    convert_reports,
    main,
)


def test_convert_report_file_splits_markdown_and_writes_static_phema(tmp_path):
    reports_dir = tmp_path / "reports"
    output_dir = tmp_path / "phemas"
    reports_dir.mkdir()

    report_path = reports_dir / "nvda_report.md"
    report_path.write_text(
        (
            "# NVDA Report\n\n"
            "AI demand is still driving top-line acceleration.\n\n"
            "## Thesis\n"
            "Revenue remains strong.\n\n"
            "## Risks\n"
            "Competition could pressure margins.\n"
        ),
        encoding="utf-8",
    )

    written_path = convert_report_file(
        report_path,
        output_dir=output_dir,
        source_root=reports_dir,
        default_owner="Attas",
        ollama_model=None,
    )

    payload = json.loads(written_path.read_text(encoding="utf-8"))

    assert written_path == output_dir / "nvda_report.json"
    assert payload["name"] == "NVDA Report"
    assert payload["owner"] == "Attas"
    assert payload["resolution_mode"] == "static"
    assert payload["meta"]["source_format"] == "markdown"
    assert payload["meta"]["source_path"] == "nvda_report.md"
    assert payload["description"] == "AI demand is still driving top-line acceleration."
    assert payload["input_schema"] == {"type": "object", "properties": {}}
    assert payload["meta"]["metadata_generated_by_ollama"] is False
    assert [section["name"] for section in payload["sections"]] == ["Introduction", "Thesis", "Risks"]
    assert payload["sections"][0]["content"] == [
        {"type": "text", "text": "AI demand is still driving top-line acceleration."}
    ]


def test_convert_report_payload_parses_dynamic_placeholders_and_input_schema():
    payload = convert_report_payload(
        {
            "name": "Company Brief",
            "body": (
                "## Snapshot\n"
                "Company: {{ company_profile.profile.name | input.symbol }}\n"
                "Price: {{ last_price.quote.price | input.symbol | pulser_id=price-pulser | param.currency=USD }}\n"
            ),
        },
        ollama_model=None,
    )

    section = payload["sections"][0]
    first_field = section["content"][1]
    second_field = section["content"][3]

    assert payload["resolution_mode"] == "dynamic"
    assert payload["input_schema"] == {
        "type": "object",
        "properties": {
            "symbol": {"type": "string"},
        },
        "required": ["symbol"],
    }
    assert section["name"] == "Snapshot"
    assert section["content"][0] == {"type": "text", "text": "Company: "}
    assert first_field["type"] == "pulse-field"
    assert first_field["pulse_name"] == "company_profile"
    assert first_field["field_path"] == "profile.name"
    assert first_field["selected_fields"] == ["profile.name"]
    assert section["content"][2] == {"type": "text", "text": "\nPrice: "}
    assert second_field["pulser_id"] == "price-pulser"
    assert second_field["params"] == {"currency": "USD"}


def test_convert_reports_handles_json_specs_and_preserves_relative_paths(tmp_path):
    reports_dir = tmp_path / "reports"
    output_dir = tmp_path / "phemas"
    nested_dir = reports_dir / "nested"
    nested_dir.mkdir(parents=True)

    report_path = nested_dir / "macro.json"
    report_path.write_text(
        json.dumps(
            {
                "title": "Macro Brief",
                "tags": ["macro", "daily"],
                "sections": [
                    {
                        "name": "Overview",
                        "content": [
                            "GDP: {{ macro_snapshot.gdp | optional.report_date=date }}"
                        ],
                    }
                ],
            }
        ),
        encoding="utf-8",
    )

    written_paths = convert_reports(
        input_dir=reports_dir,
        output_dir=output_dir,
        default_owner="MacroDesk",
        ollama_model=None,
    )

    assert written_paths == [output_dir / "nested" / "macro.json"]

    payload = json.loads((output_dir / "nested" / "macro.json").read_text(encoding="utf-8"))
    section = payload["sections"][0]

    assert payload["name"] == "Macro Brief"
    assert payload["owner"] == "MacroDesk"
    assert payload["tags"] == ["macro", "daily"]
    assert payload["meta"]["source_format"] == "json"
    assert payload["meta"]["source_path"] == "nested/macro.json"
    assert payload["resolution_mode"] == "dynamic"
    assert payload["input_schema"] == {
        "type": "object",
        "properties": {
            "report_date": {"type": "string", "format": "date"},
        },
    }
    assert section["content"][0] == {"type": "text", "text": "GDP: "}
    assert section["content"][1]["pulse_name"] == "macro_snapshot"
    assert section["content"][1]["field_path"] == "gdp"


def test_convert_report_file_reads_pdf_pages_as_sections(tmp_path):
    reports_dir = tmp_path / "reports"
    output_dir = tmp_path / "phemas"
    reports_dir.mkdir()

    report_path = reports_dir / "quarterly_review.pdf"

    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("helvetica", size=12)
    pdf.multi_cell(0, 10, "Executive Summary\nRevenue accelerated in Q1.")
    pdf.add_page()
    pdf.multi_cell(0, 10, "Risks\nCompetition remains elevated.")
    pdf.output(str(report_path))

    written_path = convert_report_file(report_path, output_dir=output_dir, source_root=reports_dir, ollama_model=None)
    payload = json.loads(written_path.read_text(encoding="utf-8"))

    assert payload["name"] == "Quarterly Review"
    assert payload["meta"]["source_format"] == "pdf"
    assert [section["name"] for section in payload["sections"]] == ["Executive Summary", "Risks"]
    assert payload["sections"][0]["content"] == [
        {"type": "text", "text": "Revenue accelerated in Q1."}
    ]
    assert "Competition remains elevated." in payload["sections"][1]["content"][0]["text"]


def test_convert_report_file_reads_pptx_slides_as_sections(tmp_path):
    reports_dir = tmp_path / "reports"
    output_dir = tmp_path / "phemas"
    reports_dir.mkdir()

    report_path = reports_dir / "board_update.pptx"

    presentation = Presentation()
    slide_one = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide_one.shapes.title.text = "Overview"
    slide_one.placeholders[1].text = "Revenue accelerated.\nMargin held steady."

    slide_two = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide_two.shapes.title.text = "Risks"
    slide_two.placeholders[1].text = "Competition increased."

    presentation.save(str(report_path))

    written_path = convert_report_file(report_path, output_dir=output_dir, source_root=reports_dir, ollama_model=None)
    payload = json.loads(written_path.read_text(encoding="utf-8"))

    assert payload["name"] == "Board Update"
    assert payload["meta"]["source_format"] == "pptx"
    assert [section["name"] for section in payload["sections"]] == ["Overview", "Risks"]
    assert payload["sections"][0]["content"] == [
        {"type": "text", "text": "Revenue accelerated.\nMargin held steady."}
    ]
    assert payload["sections"][1]["content"] == [
        {"type": "text", "text": "Competition increased."}
    ]


def test_convert_report_payload_uses_ollama_to_divide_sections_and_generate_metadata(monkeypatch):
    class FakeResponse:
        def __init__(self, payload):
            self._payload = payload
            self.status_code = 200
            self.content = json.dumps(payload).encode("utf-8")

        def raise_for_status(self):
            return None

        def json(self):
            return self._payload

    captured = {}

    def fake_post(url, json=None, timeout=None, **kwargs):
        captured["url"] = url
        captured["json"] = dict(json or {})
        captured["timeout"] = timeout
        return FakeResponse(
            {
                "response": json_module.dumps(
                    {
                        "name": "AI Infrastructure Brief",
                        "description": "A concise summary of the company profile and latest market price.",
                        "sections": [
                            {
                                "name": "Snapshot",
                                "modifier": "analytical, concise, investor-facing",
                                "description": "Key company context and price framing.",
                                "chunk_ids": [1, 2],
                            },
                            {
                                "name": "Risks",
                                "modifier": "cautious, risk-aware, investor-facing",
                                "description": "Primary downside considerations.",
                                "chunk_ids": [3],
                            }
                        ],
                    }
                )
            }
        )

    import attas.utils.report_to_phema as report_to_phema_module
    json_module = json

    monkeypatch.setattr(report_to_phema_module.requests, "post", fake_post)

    payload = convert_report_payload(
        {
            "name": "Company Brief",
            "body": (
                "Company: {{ company_profile.profile.name | input.symbol }}\n\n"
                "Price: {{ last_price.quote.price | input.symbol }}\n\n"
                "Risk: Competition remains elevated.\n"
            ),
        },
        ollama_model="custom-ollama-model",
        ollama_url=DEFAULT_OLLAMA_URL,
    )

    assert captured["url"] == DEFAULT_OLLAMA_URL
    assert captured["json"]["model"] == "custom-ollama-model"
    assert captured["json"]["keep_alive"] == DEFAULT_OLLAMA_KEEP_ALIVE
    assert captured["timeout"] == DEFAULT_OLLAMA_TIMEOUT
    assert "You must return top-level `name` and top-level `description`." in captured["json"]["prompt"]
    assert "You must return `name`, `description`, `modifier`, and `chunk_ids` for every section." in captured["json"]["prompt"]
    assert "Divide by topic, narrative, or argument, not by page or slide." in captured["json"]["prompt"]
    assert payload["name"] == "AI Infrastructure Brief"
    assert payload["description"] == "A concise summary of the company profile and latest market price."
    assert payload["sections"][0]["name"] == "Snapshot"
    assert payload["sections"][0]["modifier"] == "analytical, concise, investor-facing"
    assert payload["sections"][0]["description"] == "Key company context and price framing."
    assert payload["sections"][0]["content"][0] == {"type": "text", "text": "Company: "}
    assert payload["sections"][0]["content"][1]["pulse_name"] == "company_profile"
    assert payload["sections"][0]["content"][2] == {"type": "text", "text": "\n\nPrice: "}
    assert payload["sections"][0]["content"][3]["pulse_name"] == "last_price"
    assert payload["sections"][1]["name"] == "Risks"
    assert payload["sections"][1]["modifier"] == "cautious, risk-aware, investor-facing"
    assert payload["sections"][1]["description"] == "Primary downside considerations."
    assert payload["sections"][1]["content"] == [
        {"type": "text", "text": "Risk: Competition remains elevated."}
    ]
    assert payload["meta"]["metadata_generated_by_ollama"] is True
    assert payload["meta"]["metadata_model"] == "custom-ollama-model"


def test_convert_report_file_appends_model_name_to_output_filename(tmp_path):
    reports_dir = tmp_path / "reports"
    output_dir = tmp_path / "phemas"
    reports_dir.mkdir()

    report_path = reports_dir / "simple_report.md"
    report_path.write_text("# Simple Report\n\nRevenue improved.\n", encoding="utf-8")

    written_path = convert_report_file(
        report_path,
        output_dir=output_dir,
        source_root=reports_dir,
        ollama_model="qwen3:8b",
    )

    assert written_path == output_dir / "simple_report_qwen3_8b.json"


def test_main_emits_detail_logs_and_stdout_summary(tmp_path, capsys):
    reports_dir = tmp_path / "reports"
    output_dir = tmp_path / "phemas"
    reports_dir.mkdir()

    report_path = reports_dir / "simple_report.md"
    report_path.write_text(
        "# Simple Report\n\n## Overview\nRevenue improved.\n",
        encoding="utf-8",
    )

    exit_code = main(
        [
            "--input-dir",
            str(reports_dir),
            "--output-dir",
            str(output_dir),
            "--model",
            "",
            "--verbose",
        ]
    )

    captured = capsys.readouterr()
    summary_csv_path = output_dir / "model_performance_summary.csv"

    assert exit_code == 0
    assert "[INFO] Starting batch conversion" in captured.out
    assert "[DEBUG] Prepared" in captured.out
    assert "[INFO] Converted report:" in captured.out
    assert "[INFO] Selected 1 model(s): <disabled>" in captured.out
    assert "Converted 1 report(s) into" in captured.out
    assert str((output_dir / "simple_report.json").resolve()) in captured.out
    assert f"Model performance CSV: {summary_csv_path.resolve()}" in captured.out
    with summary_csv_path.open("r", encoding="utf-8", newline="") as handle:
        rows = list(csv.DictReader(handle))
    assert rows == [
        {
            "model_name": "<disabled>",
            "filename": "simple_report.md",
            "number_of_sections": "1",
            "number_of_names": "1",
            "number_of_descriptions": "0",
            "number_of_modifiers": "0",
            "process_time_seconds": rows[0]["process_time_seconds"],
        }
    ]
    assert float(rows[0]["process_time_seconds"]) >= 0.0


def test_main_uses_all_default_models_when_model_flag_is_omitted(monkeypatch, tmp_path, capsys):
    import attas.utils.report_to_phema as report_to_phema_module

    events = []

    def fake_convert_reports(**kwargs):
        events.append(("convert", kwargs["ollama_model"], kwargs["ollama_keep_alive"]))
        output_dir = Path(kwargs["output_dir"])
        model = kwargs["ollama_model"]
        suffix = model.replace(":", "_").replace(".", "_")
        path = output_dir / f"report_{suffix}.json"
        return [path]

    def fake_warm_ollama_model(model, **kwargs):
        events.append(("warm", model, kwargs["keep_alive"]))
        return True

    monkeypatch.setattr(report_to_phema_module, "convert_reports", fake_convert_reports)
    monkeypatch.setattr(report_to_phema_module, "_warm_ollama_model", fake_warm_ollama_model)

    exit_code = main(
        [
            "--input-dir",
            str(tmp_path / "reports"),
            "--output-dir",
            str(tmp_path / "phemas"),
        ]
    )

    captured = capsys.readouterr()

    assert exit_code == 0
    assert [event[1] for event in events if event[0] == "warm"] == DEFAULT_OLLAMA_MODELS
    assert [event[1] for event in events if event[0] == "convert"] == DEFAULT_OLLAMA_MODELS
    assert all(event[2] == DEFAULT_OLLAMA_KEEP_ALIVE for event in events)
    for model in DEFAULT_OLLAMA_MODELS:
        assert events.index(("warm", model, DEFAULT_OLLAMA_KEEP_ALIVE)) < events.index(
            ("convert", model, DEFAULT_OLLAMA_KEEP_ALIVE)
        )
    assert f"[INFO] Selected {len(DEFAULT_OLLAMA_MODELS)} model(s):" in captured.out
    assert f"Converted {len(DEFAULT_OLLAMA_MODELS)} report(s) into" in captured.out


def test_main_refines_prompt_template_and_saves_best_candidate(monkeypatch, tmp_path, capsys):
    import attas.utils.report_to_phema as report_to_phema_module

    reports_dir = tmp_path / "reports"
    reports_dir.mkdir()
    report_path = reports_dir / "sample_report.md"
    report_path.write_text("# Sample Report\n\nRevenue improved.\n", encoding="utf-8")

    prompt_template_path = tmp_path / "section_plan_prompt_template.txt"
    original_template = (
        "base prompt\n"
        "[[CURRENT_NAME_HINT]]\n"
        "[[CURRENT_DESCRIPTION_HINT]]\n"
        "[[CHUNK_BLOCK]]\n"
    )
    prompt_template_path.write_text(original_template, encoding="utf-8")

    candidates = iter(
        [
            "candidate one\n[[CURRENT_NAME_HINT]]\n[[CURRENT_DESCRIPTION_HINT]]\n[[CHUNK_BLOCK]]\n",
            "candidate two\n[[CURRENT_NAME_HINT]]\n[[CURRENT_DESCRIPTION_HINT]]\n[[CHUNK_BLOCK]]\n",
        ]
    )

    def fake_warm_ollama_model(*args, **kwargs):
        return True

    def fake_generate_refined_prompt_template(**kwargs):
        return next(candidates)

    def fake_evaluate_prompt_template(**kwargs):
        template = kwargs["prompt_template_text"]
        if template.startswith("candidate two"):
            payload = {
                "name": "Improved Report",
                "description": "Better prompt output.",
                "sections": [
                    {"name": "Overview", "description": "Overview section.", "modifier": "analytical"},
                    {"name": "Drivers", "description": "Drivers section.", "modifier": "data-rich"},
                    {"name": "Risks", "description": "Risk section.", "modifier": "cautious"},
                ],
            }
            score = (2, 3, 3, 3, 3, 80)
        elif template.startswith("candidate one"):
            payload = {
                "name": "Candidate Report",
                "description": "Not much better.",
                "sections": [
                    {"name": "Overview", "description": "", "modifier": ""},
                ],
            }
            score = (2, 0, 1, 0, 0, 20)
        else:
            payload = {
                "name": "Baseline Report",
                "description": "Baseline prompt output.",
                "sections": [
                    {"name": "Overview", "description": "Overview section.", "modifier": ""},
                ],
            }
            score = (2, 0, 1, 1, 0, 30)
        return {"payload": payload, "score": score}

    monkeypatch.setattr(report_to_phema_module, "_warm_ollama_model", fake_warm_ollama_model)
    monkeypatch.setattr(report_to_phema_module, "_generate_refined_prompt_template", fake_generate_refined_prompt_template)
    monkeypatch.setattr(report_to_phema_module, "_evaluate_prompt_template", fake_evaluate_prompt_template)

    exit_code = main(
        [
            "--input-dir",
            str(reports_dir),
            "--prompt-template",
            str(prompt_template_path),
            "--refine-prompt",
            "qwen3:8b",
            "--file",
            report_path.name,
            "--times",
            "2",
        ]
    )

    captured = capsys.readouterr()
    backup_files = list(tmp_path.glob("section_plan_prompt_template.*.bak.txt"))

    assert exit_code == 0
    assert prompt_template_path.read_text(encoding="utf-8").startswith("candidate two")
    assert len(backup_files) == 1
    assert backup_files[0].read_text(encoding="utf-8") == original_template
    assert "Saved best prompt template to" in captured.out
    assert "Best prompt produced 3 section(s)" in captured.out


def test_main_refine_prompt_stops_early_when_all_sections_have_metadata(monkeypatch, tmp_path, capsys):
    import attas.utils.report_to_phema as report_to_phema_module

    reports_dir = tmp_path / "reports"
    reports_dir.mkdir()
    report_path = reports_dir / "sample_report.md"
    report_path.write_text("# Sample Report\n\nRevenue improved.\n", encoding="utf-8")

    prompt_template_path = tmp_path / "section_plan_prompt_template.txt"
    original_template = (
        "base prompt\n"
        "[[CURRENT_NAME_HINT]]\n"
        "[[CURRENT_DESCRIPTION_HINT]]\n"
        "[[CHUNK_BLOCK]]\n"
    )
    prompt_template_path.write_text(original_template, encoding="utf-8")

    calls = {"generated": 0}

    def fake_warm_ollama_model(*args, **kwargs):
        return True

    def fake_generate_refined_prompt_template(**kwargs):
        calls["generated"] += 1
        return "candidate one\n[[CURRENT_NAME_HINT]]\n[[CURRENT_DESCRIPTION_HINT]]\n[[CHUNK_BLOCK]]\n"

    def fake_evaluate_prompt_template(**kwargs):
        template = kwargs["prompt_template_text"]
        if template.startswith("candidate one"):
            payload = {
                "name": "Improved Report",
                "description": "Improved prompt output.",
                "sections": [
                    {"name": "Overview", "description": "Overview section.", "modifier": "analytical"},
                    {"name": "Drivers", "description": "Drivers section.", "modifier": "data-rich"},
                ],
            }
            score = (2, 2, 2, 2, 2, 70)
        else:
            payload = {
                "name": "Baseline Report",
                "description": "Baseline prompt output.",
                "sections": [
                    {"name": "Overview", "description": "Overview section.", "modifier": ""},
                ],
            }
            score = (2, 0, 1, 1, 0, 20)
        return {"payload": payload, "score": score}

    monkeypatch.setattr(report_to_phema_module, "_warm_ollama_model", fake_warm_ollama_model)
    monkeypatch.setattr(report_to_phema_module, "_generate_refined_prompt_template", fake_generate_refined_prompt_template)
    monkeypatch.setattr(report_to_phema_module, "_evaluate_prompt_template", fake_evaluate_prompt_template)

    exit_code = main(
        [
            "--input-dir",
            str(reports_dir),
            "--prompt-template",
            str(prompt_template_path),
            "--refine-prompt",
            "qwen3:8b",
            "--file",
            report_path.name,
            "--times",
            "5",
        ]
    )

    captured = capsys.readouterr()

    assert exit_code == 0
    assert calls["generated"] == 1
    assert prompt_template_path.read_text(encoding="utf-8").startswith("candidate one")
    assert "Stopping prompt refinement early because all sections already contain name, description, and modifier" in captured.out
