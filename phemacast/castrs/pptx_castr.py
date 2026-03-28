from __future__ import annotations

import json
import os
import uuid
from datetime import datetime, timezone
from typing import Any, Dict, List, Optional

from phemacast.agents.castr import Castr

try:  # pragma: no cover - exercised in integration tests when dependency exists
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
except ImportError:  # pragma: no cover - keeps module importable when optional dep is absent
    Presentation = None
    RGBColor = None
    PP_ALIGN = None
    Inches = None
    Pt = None


def _normalize_text(value: Any) -> str:
    if value is None:
        return ""
    if isinstance(value, str):
        return value.strip()
    if isinstance(value, (int, float, bool)):
        return str(value)
    if isinstance(value, list):
        parts = [_normalize_text(item) for item in value]
        return "; ".join(part for part in parts if part)
    if isinstance(value, dict):
        for key in ("rendered", "value", "content", "description", "summary", "text"):
            text = _normalize_text(value.get(key))
            if text:
                return text
        if "result" in value:
            label = _normalize_text(value.get("name") or value.get("key") or value.get("pulse_name"))
            text = _normalize_text(value.get("result"))
            if label and text:
                return f"{label}: {text}"
            if text:
                return text
        if "data" in value:
            label = _normalize_text(value.get("name") or value.get("key") or value.get("pulse_name"))
            text = _normalize_text(value.get("data"))
            if label and text:
                return f"{label}: {text}"
            if text:
                return text
        compact = json.dumps(value, ensure_ascii=True, sort_keys=True, default=str)
        return compact
    return str(value).strip()


def _section_lines(section: Dict[str, Any]) -> List[str]:
    lines: List[str] = []

    description = _normalize_text(section.get("description"))
    if description:
        lines.append(description)

    modifier = _normalize_text(section.get("modifier"))
    if modifier:
        lines.append(f"Modifier: {modifier}")

    content = section.get("content", [])
    if isinstance(content, list):
        for item in content:
            text = _normalize_text(item)
            if not text:
                continue
            lines.extend(part.strip() for part in text.splitlines() if part.strip())
    else:
        text = _normalize_text(content)
        if text:
            lines.extend(part.strip() for part in text.splitlines() if part.strip())

    return lines or ["No section content provided."]


def _chunk_lines(lines: List[str], size: int = 6) -> List[List[str]]:
    if not lines:
        return [["No section content provided."]]
    return [lines[index:index + size] for index in range(0, len(lines), size)]


class PPTXCastr(Castr):
    """Castr implementation that renders a Phema into a real PowerPoint deck."""

    def cast(
        self,
        phema: Dict[str, Any],
        format: str,
        preferences: Optional[Dict[str, Any]] = None,
    ) -> Dict[str, Any]:
        target_format = "PPTX"
        media_id = str(uuid.uuid4())
        filename = f"{media_id}.pptx"

        if not self.pool or not hasattr(self.pool, "root_path"):
            raise ValueError("FileSystemPool with root_path is required for PPTXCastr.")

        media_dir = os.path.join(self.pool.root_path, "media")
        os.makedirs(media_dir, exist_ok=True)
        file_path = os.path.join(media_dir, filename)

        self._generate_real_pptx(file_path, phema, preferences)

        return {
            "status": "success",
            "media_id": media_id,
            "format": target_format,
            "message": f"Successfully generated PPTX for Phema: {phema.get('name', 'Untitled')}",
            "location": f"media/{filename}",
            "url": f"/api/media/{filename}",
        }

    def _generate_real_pptx(
        self,
        file_path: str,
        phema: Dict[str, Any],
        preferences: Optional[Dict[str, Any]] = None,
    ) -> None:
        if Presentation is None:
            raise RuntimeError("python-pptx is required to generate PPTX output.")

        pref = preferences or {}
        audience = _normalize_text(pref.get("audience")) or "General"
        language = _normalize_text(pref.get("language")) or "en"
        theme = _normalize_text(pref.get("theme")) or "Executive"
        generated_at = datetime.now(timezone.utc).strftime("%Y-%m-%d %H:%M:%S UTC")

        title = _normalize_text(phema.get("name")) or "Untitled Phema"
        description = _normalize_text(phema.get("description")) or "No description provided."
        sections = phema.get("sections", [])
        if not isinstance(sections, list):
            sections = []

        presentation = Presentation()
        presentation.slide_width = Inches(13.333)
        presentation.slide_height = Inches(7.5)

        core = presentation.core_properties
        core.title = title
        core.subject = description
        core.author = self.name
        core.keywords = "phemacast, castr, pptx"
        core.comments = "Generated by Phemacast PPTXCastr."

        self._add_title_slide(
            presentation,
            title=title,
            description=description,
            audience=audience,
            generated_at=generated_at,
        )
        self._add_metadata_slide(
            presentation,
            phema=phema,
            audience=audience,
            language=language,
            theme=theme,
            generated_at=generated_at,
        )

        for index, raw_section in enumerate(sections, start=1):
            section = raw_section if isinstance(raw_section, dict) else {"name": f"Section {index}", "content": [raw_section]}
            section_name = _normalize_text(section.get("name")) or f"Section {index}"
            line_groups = _chunk_lines(_section_lines(section))
            for group_index, lines in enumerate(line_groups, start=1):
                slide_title = section_name if group_index == 1 else f"{section_name} (cont.)"
                self._add_section_slide(presentation, slide_title, lines)

        if len(presentation.slides) == 2:
            self._add_section_slide(presentation, "Summary", [description])

        presentation.save(file_path)

    def _add_title_slide(
        self,
        presentation: "Presentation",
        *,
        title: str,
        description: str,
        audience: str,
        generated_at: str,
    ) -> None:
        slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        self._set_background(slide, RGBColor(25, 45, 70))

        title_shape = slide.shapes.title
        title_shape.text = title
        title_paragraph = title_shape.text_frame.paragraphs[0]
        title_paragraph.alignment = PP_ALIGN.LEFT
        title_run = title_paragraph.runs[0]
        title_run.font.size = Pt(28)
        title_run.font.bold = True
        title_run.font.color.rgb = RGBColor(255, 255, 255)

        subtitle = slide.placeholders[1]
        subtitle.text = (
            f"{description}\n\n"
            f"Audience: {audience}\n"
            f"Generated: {generated_at}"
        )
        for paragraph in subtitle.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(16)
                run.font.color.rgb = RGBColor(225, 233, 242)

    def _add_metadata_slide(
        self,
        presentation: "Presentation",
        *,
        phema: Dict[str, Any],
        audience: str,
        language: str,
        theme: str,
        generated_at: str,
    ) -> None:
        tags = phema.get("tags", [])
        tag_text = ", ".join(str(tag) for tag in tags if str(tag).strip()) or "None"
        owner = _normalize_text(phema.get("owner")) or "Unspecified"
        phema_id = _normalize_text(phema.get("phema_id") or phema.get("id")) or "N/A"
        lines = [
            f"Owner: {owner}",
            f"Phema ID: {phema_id}",
            f"Target audience: {audience}",
            f"Language: {language}",
            f"Theme: {theme}",
            f"Tags: {tag_text}",
            f"Generated: {generated_at}",
        ]
        self._add_section_slide(presentation, "Document Metadata", lines)

    def _add_section_slide(
        self,
        presentation: "Presentation",
        title: str,
        lines: List[str],
    ) -> None:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        title_shape = slide.shapes.title
        title_shape.text = title
        title_paragraph = title_shape.text_frame.paragraphs[0]
        title_paragraph.alignment = PP_ALIGN.LEFT
        title_run = title_paragraph.runs[0]
        title_run.font.size = Pt(24)
        title_run.font.bold = True
        title_run.font.color.rgb = RGBColor(33, 56, 82)

        body_shape = self._resolve_body_shape(slide)
        text_frame = body_shape.text_frame
        text_frame.clear()
        text_frame.word_wrap = True

        for index, line in enumerate(lines):
            paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
            paragraph.text = line
            paragraph.level = 0
            paragraph.space_after = Pt(10)
            for run in paragraph.runs:
                run.font.size = Pt(18)
                run.font.color.rgb = RGBColor(45, 45, 45)

    def _resolve_body_shape(self, slide: Any) -> Any:
        try:
            return slide.placeholders[1]
        except (IndexError, KeyError):
            return slide.shapes.add_textbox(Inches(0.9), Inches(1.7), Inches(11.5), Inches(4.8))

    def _set_background(self, slide: Any, color: "RGBColor") -> None:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color


PptxCastr = PPTXCastr
PptCastr = PPTXCastr

__all__ = ["PPTXCastr", "PptxCastr", "PptCastr"]
